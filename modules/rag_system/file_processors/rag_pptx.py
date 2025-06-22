import logging

def extract_text(file_path: str, **kwargs) -> dict:
    result = {
        "text": "",
        "success": False,
        "error": None,
        "cleaned": True,
        "meta": {
            "file_path": file_path,
            "file_type": "pptx",
            "parser": "rag_pptx",
            "lines": 0,
            "chars": 0,
        }
    }
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slides_text.append(shape.text.strip())
        cleaned_text = "\n".join([t for t in slides_text if t])
        result["text"] = cleaned_text
        result["success"] = True
        result["meta"]["lines"] = cleaned_text.count('\n') + 1
        result["meta"]["chars"] = len(cleaned_text)
    except Exception as e:
        result["error"] = str(e)
        logging.getLogger("rag_pptx").warning(f"PPTX extraction failed for {file_path}: {e}")
    return result
